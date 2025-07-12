import io
from fastapi import UploadFile, HTTPException

import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from ..config import logger

def _extract_text_from_pdf(file_stream: io.BytesIO) -> str:
    """Extracts text from a PDF file stream."""
    text = ""
    with fitz.open(stream=file_stream, filetype="pdf") as doc:
        for page in doc:
            text += page.get_text()
    return text

def _extract_text_from_docx(file_stream: io.BytesIO) -> str:
    """Extracts text from a DOCX file stream."""
    document = Document(file_stream)
    return "\n".join([para.text for para in document.paragraphs])

def _extract_text_from_pptx(file_stream: io.BytesIO) -> str:
    """Extracts text from a PPTX file stream."""
    prs = Presentation(file_stream)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def _extract_text_from_txt(file_stream: io.BytesIO) -> str:
    """Extracts text from a TXT or MD file stream."""
    return file_stream.getvalue().decode("utf-8")

async def extract_text_from_file(file: UploadFile) -> str:
    """
    Reads an uploaded file and dispatches to the correct text extraction function
    based on the file's content type or extension.
    """
    file_stream = io.BytesIO(await file.read())
    filename = file.filename.lower()

    if filename.endswith(".pdf"):
        return _extract_text_from_pdf(file_stream)
    elif filename.endswith(".docx"):
        return _extract_text_from_docx(file_stream)
    elif filename.endswith(".pptx"):
        return _extract_text_from_pptx(file_stream)
    elif filename.endswith((".txt", ".md")):
        return _extract_text_from_txt(file_stream)
    else:
        raise HTTPException(status_code=400, detail=f"Unsupported file type: {file.filename}")

from . import ai_service
from pydantic import BaseModel, Field
from typing import List

class ParsedItems(BaseModel):
    items: List[str] = Field(default_factory=list)

AI_PARSER_PROMPT = """
You are a highly intelligent text analysis assistant for an application called WordNest. Your primary function is to read a large piece of text and break it down into a list of valuable, atomic learning items for the user. These items can be individual words, phrases, or full sentences.

**Your Task:**
Analyze the user-provided text and extract a list of all noteworthy items.

**Rules:**
1.  **Extract Everything Valuable:** Your goal is to be comprehensive. Identify and pull out:
    *   **Key Vocabulary:** Important or technical single words.
    *   **Important Phrases:** Idioms, collocations, or technical terms composed of multiple words (e.g., "supply chain management", "return on investment").
    *   **Well-Constructed Sentences:** Sentences that are great examples of grammar, style, or that contain a key idea.
2.  **Maintain Original Form:** The extracted items in the list must be the exact strings as they appeared in the original text.
3.  **Return JSON format:** Your final output **MUST** be a JSON object with a single key, `items`, which holds a list of the extracted strings. If no items are found, return an empty list: `{"items": []}`.
4.  **Be thorough:** Do not summarize. Extract directly. The user will review and select which items to save.

**Example Input Text:**
"The new manufacturing process, which relies on just-in-time inventory management, has significantly reduced overhead costs."

**Example JSON Output:**
{
  "items": [
    "manufacturing process",
    "just-in-time inventory management",
    "significantly reduced",
    "overhead costs"
  ]
}

**Now, analyze the following text and provide the JSON output:**
"""

def extract_learning_items_from_text(text: str) -> dict:
    """
    Sends the extracted text to an AI service to get a list of learning items (words, phrases, sentences).
    """
    json_response = ai_service.call_ai(
        system_prompt=AI_PARSER_PROMPT,
        user_prompt=text
    )

    if not json_response:
        return {"items": []}
    
    try:
        validated_response = ParsedItems.model_validate(json_response)
        return validated_response.model_dump()
    except Exception as e:
        logger.error(f"Failed to validate parser AI response: {e}")
        return {"items": []}